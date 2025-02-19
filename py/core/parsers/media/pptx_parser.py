# type: ignore
from io import BytesIO
from typing import AsyncGenerator

from pptx import Presentation

from core.base.parsers.base_parser import AsyncParser
from core.base.providers import (
    CompletionProvider,
    DatabaseProvider,
    IngestionConfig,
)


class PPTXParser(AsyncParser[str | bytes]):
    """A parser for PPT data."""

    def __init__(
        self,
        config: IngestionConfig,
        database_provider: DatabaseProvider,
        llm_provider: CompletionProvider,
    ):
        self.database_provider = database_provider
        self.llm_provider = llm_provider
        self.config = config
        self.Presentation = Presentation

    async def ingest(
        self, data: str | bytes, **kwargs
    ) -> AsyncGenerator[str, None]:  # type: ignore
        """Ingest PPT data and yield text from each slide."""
        if isinstance(data, str):
            raise ValueError("PPT data must be in bytes format.")

        prs = self.Presentation(BytesIO(data))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    yield shape.text
