import json
import string
from abc import ABC, abstractmethod
from io import BytesIO
from typing import Any, Generator, Generic, TypeVar, Union

from bs4 import BeautifulSoup

T = TypeVar("T")


class Ingestor(ABC, Generic[T]):
    @abstractmethod
    def ingest(self, data: Any) -> Generator[str, None, None]:
        pass


class TextAdapter(Ingestor[str]):
    def ingest(self, data: str) -> Generator[str, None, None]:
        yield data


class JSONAdapter(Ingestor[str]):
    def ingest(self, data: Union[str, bytes]) -> Generator[str, None, None]:
        if isinstance(data, bytes):
            data = data.decode("utf-8")
        yield self._parse_json(json.loads(data))

    def _parse_json(self, data: dict) -> str:
        def remove_objects_with_null(obj):
            if not isinstance(obj, dict):
                return obj
            result = obj.copy()
            for key, value in obj.items():
                if isinstance(value, dict):
                    result[key] = remove_objects_with_null(value)
                elif value is None:
                    del result[key]
            return result

        def format_json_as_text(obj, indent=0):
            lines = []
            indent_str = " " * indent

            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, (dict, list)):
                        nested = format_json_as_text(value, indent + 2)
                        lines.append(f"{indent_str}{key}:\n{nested}")
                    else:
                        lines.append(f"{indent_str}{key}: {value}")
            elif isinstance(obj, list):
                for item in obj:
                    nested = format_json_as_text(item, indent + 2)
                    lines.append(f"{nested}")
            else:
                return f"{indent_str}{obj}"

            return "\n".join(lines)

        return format_json_as_text(remove_objects_with_null(data))


class HTMLAdapter(Ingestor[str]):
    def ingest(self, data: str) -> list[str]:
        soup = BeautifulSoup(data, "html.parser")
        yield soup.get_text()


class PDFAdapter(Ingestor[str]):
    def __init__(self):
        try:
            from pypdf import PdfReader

            self.PdfReader = PdfReader
        except ImportError:
            raise ValueError(
                "Error, `pypdf` is required to run `PyPDFAdapter`. Please install it using `pip install pypdf`."
            )

    def ingest(self, data: bytes) -> Generator[str, None, None]:
        pdf = self.PdfReader(BytesIO(data))
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text is not None:
                page_text = "".join(
                    filter(lambda x: x in string.printable, page_text)
                )
                yield page_text


class PPTAdapter(Ingestor[str]):
    def __init__(self):
        try:
            from pptx import Presentation

            self.Presentation = Presentation
        except ImportError:
            raise ValueError(
                "Error, `python-pptx` is required to run `PPTAdapter`. Please install it using `pip install python-pptx`."
            )

    def ingest(self, data: bytes) -> list[str]:
        prs = self.Presentation(BytesIO(data))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    yield shape.text


class DOCXAdapter(Ingestor[str]):
    def __init__(self):
        try:
            from docx import Document

            self.Document = Document
        except ImportError:
            raise ValueError(
                "Error, `python-docx` is required to run `DOCXAdapter`. Please install it using `pip install python-docx`."
            )

    def ingest(self, data: bytes) -> Generator[str, None, None]:
        doc = self.Document(BytesIO(data))
        for paragraph in doc.paragraphs:
            yield paragraph.text


class CSVAdapter(Ingestor[str]):
    def __init__(self):
        import csv
        from io import StringIO

        self.csv = csv
        self.StringIO = StringIO

    def ingest(self, data: Union[str, bytes]) -> Generator[str, None, None]:
        if isinstance(data, bytes):
            data = data.decode("utf-8")
        csv_reader = self.csv.reader(self.StringIO(data))
        for row in csv_reader:
            yield ", ".join(row)


class XLSXAdapter(Ingestor[str]):
    def __init__(self):
        try:
            from openpyxl import load_workbook

            self.load_workbook = load_workbook
        except ImportError:
            raise ValueError(
                "Error, `openpyxl` is required to run `XLSXAdapter`. Please install it using `pip install openpyxl`."
            )

    def ingest(self, data: bytes) -> Generator[str, None, None]:
        wb = self.load_workbook(filename=BytesIO(data))
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                yield ", ".join(map(str, row))


class MarkdownAdapter(Ingestor[str]):
    def __init__(self):
        import markdown

        self.markdown = markdown

    def ingest(self, data: str) -> Generator[str, None, None]:
        html = self.markdown.markdown(data)
        soup = BeautifulSoup(html, "html.parser")
        yield soup.get_text()
