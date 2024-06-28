from io import BytesIO
from typing import AsyncGenerator

from r2r.base.abstractions.document import DataType
from r2r.base.parsers.base_parser import AsyncParser


class PPTParser(AsyncParser[DataType]):
    """A parser for PPT data."""

    def __init__(self):
        try:
            from pptx import Presentation

            self.Presentation = Presentation
        except ImportError:
            raise ValueError(
                "Error, `python-pptx` is required to run `PPTParser`. Please install it using `pip install python-pptx`."
            )

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest PPT data and yield text from each slide."""
        if isinstance(data, str):
            raise ValueError("PPT data must be in bytes format.")

        prs = self.Presentation(BytesIO(data))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    yield shape.text
