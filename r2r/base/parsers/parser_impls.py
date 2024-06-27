"""Implementations of parsers for different data types."""

import base64
import json
import os
import string
from io import BytesIO
from typing import AsyncGenerator, Union

import requests
from bs4 import BeautifulSoup

from ..abstractions.document import DataType
from .base_parser import AsyncParser


class TextParser(AsyncParser[DataType]):
    """A parser for raw text data."""

    async def ingest(self, data: DataType) -> AsyncGenerator[DataType, None]:
        if isinstance(data, bytes):
            data = data.decode("utf-8")
        yield data


class JSONParser(AsyncParser[DataType]):
    """A parser for JSON data."""

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest JSON data and yield a formatted text representation."""
        if isinstance(data, bytes):
            data = data.decode("utf-8")
        yield self._parse_json(json.loads(data))

    def _parse_json(self, data: dict) -> str:
        def remove_objects_with_null(obj):
            if not isinstance(obj, dict):
                return obj
            result = obj.copy()
            for key, value in obj.items():
                if isinstance(value, dict):
                    result[key] = remove_objects_with_null(value)
                elif value is None:
                    del result[key]
            return result

        def format_json_as_text(obj, indent=0):
            lines = []
            indent_str = " " * indent

            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, (dict, list)):
                        nested = format_json_as_text(value, indent + 2)
                        lines.append(f"{indent_str}{key}:\n{nested}")
                    else:
                        lines.append(f"{indent_str}{key}: {value}")
            elif isinstance(obj, list):
                for item in obj:
                    nested = format_json_as_text(item, indent + 2)
                    lines.append(f"{nested}")
            else:
                return f"{indent_str}{obj}"

            return "\n".join(lines)

        return format_json_as_text(remove_objects_with_null(data))


class HTMLParser(AsyncParser[DataType]):
    """A parser for HTML data."""

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest HTML data and yield text."""
        soup = BeautifulSoup(data, "html.parser")
        yield soup.get_text()


class PDFParser(AsyncParser[DataType]):
    """A parser for PDF data."""

    def __init__(self):
        try:
            from pypdf import PdfReader

            self.PdfReader = PdfReader
        except ImportError:
            raise ValueError(
                "Error, `pypdf` is required to run `PyPDFParser`. Please install it using `pip install pypdf`."
            )

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest PDF data and yield text from each page."""
        if isinstance(data, str):
            raise ValueError("PDF data must be in bytes format.")

        pdf = self.PdfReader(BytesIO(data))
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text is not None:
                page_text = "".join(
                    filter(lambda x: x in string.printable, page_text)
                )
                yield page_text


class PPTParser(AsyncParser[DataType]):
    """A parser for PPT data."""

    def __init__(self):
        try:
            from pptx import Presentation

            self.Presentation = Presentation
        except ImportError:
            raise ValueError(
                "Error, `python-pptx` is required to run `PPTParser`. Please install it using `pip install python-pptx`."
            )

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest PPT data and yield text from each slide."""
        if isinstance(data, str):
            raise ValueError("PPT data must be in bytes format.")

        prs = self.Presentation(BytesIO(data))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    yield shape.text


class DOCXParser(AsyncParser[DataType]):
    """A parser for DOCX data."""

    def __init__(self):
        try:
            from docx import Document

            self.Document = Document
        except ImportError:
            raise ValueError(
                "Error, `python-docx` is required to run `DOCXParser`. Please install it using `pip install python-docx`."
            )

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest DOCX data and yield text from each paragraph."""
        if isinstance(data, str):
            raise ValueError("DOCX data must be in bytes format.")

        doc = self.Document(BytesIO(data))
        for paragraph in doc.paragraphs:
            yield paragraph.text


class XLSXParser(AsyncParser[DataType]):
    """A parser for XLSX data."""

    def __init__(self):
        try:
            from openpyxl import load_workbook

            self.load_workbook = load_workbook
        except ImportError:
            raise ValueError(
                "Error, `openpyxl` is required to run `XLSXParser`. Please install it using `pip install openpyxl`."
            )

    async def ingest(self, data: bytes) -> AsyncGenerator[str, None]:
        """Ingest XLSX data and yield text from each row."""
        if isinstance(data, str):
            raise ValueError("XLSX data must be in bytes format.")

        wb = self.load_workbook(filename=BytesIO(data))
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                yield ", ".join(map(str, row))


class CSVParser(AsyncParser[DataType]):
    """A parser for CSV data."""

    def __init__(self):
        import csv
        from io import StringIO

        self.csv = csv
        self.StringIO = StringIO

    async def ingest(
        self, data: Union[str, bytes]
    ) -> AsyncGenerator[str, None]:
        """Ingest CSV data and yield text from each row."""
        if isinstance(data, bytes):
            data = data.decode("utf-8")
        csv_reader = self.csv.reader(self.StringIO(data))
        for row in csv_reader:
            yield ", ".join(row)


class MarkdownParser(AsyncParser[DataType]):
    """A parser for Markdown data."""

    def __init__(self):
        import markdown

        self.markdown = markdown

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest Markdown data and yield text."""
        if isinstance(data, bytes):
            data = data.decode("utf-8")
        html = self.markdown.markdown(data)
        soup = BeautifulSoup(html, "html.parser")
        yield soup.get_text()


def process_frame_with_openai(
    data: bytes,
    api_key: str,
    model: str = "gpt-4o",
    max_tokens: int = 2_048,
    api_base: str = "https://api.openai.com/v1/chat/completions",
) -> str:
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }

    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "First, provide a title for the image, then explain everything that you see. Be very thorough in your analysis as a user will need to understand the image without seeing it. If it is possible to transcribe the image to text directly, then do so. The more detail you provide, the better the user will understand the image.",
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{data}"},
                    },
                ],
            }
        ],
        "max_tokens": max_tokens,
    }

    response = requests.post(api_base, headers=headers, json=payload)
    response_json = response.json()
    return response_json["choices"][0]["message"]["content"]


def process_audio_with_openai(
    audio_file,
    api_key: str,
    audio_api_base: str = "https://api.openai.com/v1/audio/transcriptions",
) -> str:
    headers = {"Authorization": f"Bearer {api_key}"}

    transcription_response = requests.post(
        audio_api_base,
        headers=headers,
        files={"file": audio_file},
        data={"model": "whisper-1"},
    )
    transcription = transcription_response.json()

    return transcription["text"]


class AudioParser(AsyncParser[bytes]):
    """A parser for audio data."""

    def __init__(
        self, api_base: str = "https://api.openai.com/v1/audio/transcriptions"
    ):
        self.api_base = api_base
        self.openai_api_key = os.environ.get("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError(
                "Error, environment variable `OPENAI_API_KEY` is required to run `AudioParser`."
            )

    async def ingest(self, data: bytes) -> AsyncGenerator[str, None]:
        """Ingest audio data and yield a transcription."""
        temp_audio_path = "temp_audio.wav"
        with open(temp_audio_path, "wb") as f:
            f.write(data)
        try:
            transcription_text = process_audio_with_openai(
                open(temp_audio_path, "rb"), self.openai_api_key
            )
            yield transcription_text
        finally:
            os.remove(temp_audio_path)


class ImageParser(AsyncParser[DataType]):
    """A parser for image data."""

    def __init__(
        self,
        model: str = "gpt-4o",
        max_tokens: int = 2_048,
        api_base: str = "https://api.openai.com/v1/chat/completions",
    ):
        self.model = model
        self.max_tokens = max_tokens
        self.openai_api_key = os.environ.get("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError(
                "Error, environment variable `OPENAI_API_KEY` is required to run `ImageParser`."
            )
        self.api_base = api_base

    async def ingest(self, data: DataType) -> AsyncGenerator[str, None]:
        """Ingest image data and yield a description."""
        if isinstance(data, bytes):
            import base64

            data = base64.b64encode(data).decode("utf-8")

        yield process_frame_with_openai(
            data,
            self.openai_api_key,
            self.model,
            self.max_tokens,
            self.api_base,
        )


class MovieParser(AsyncParser):
    """A parser for movie data."""

    def __init__(
        self,
        model: str = "gpt-4o",
        max_tokens: int = 2048,
        seconds_per_frame: int = 2,
        max_frames: int = 10,
    ):
        try:
            import cv2

            self.cv2 = cv2
        except ImportError:
            raise ValueError(
                "Error, `opencv-python` is required to run `MovieParser`. Please install it using `pip install opencv-python`."
            )
        try:
            import moviepy.editor as mp

            self.mp = mp
        except ImportError:
            raise ValueError(
                "Error, `moviepy` is required to run `MovieParser`. Please install it using `pip install moviepy`."
            )

        self.model = model
        self.max_tokens = max_tokens
        self.seconds_per_frame = seconds_per_frame
        self.max_frames = max_frames
        self.openai_api_key = os.environ.get("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError(
                "Error, environment variable `OPENAI_API_KEY` is required to run `MovieParser`."
            )

    async def ingest(self, data: bytes) -> AsyncGenerator[str, None]:
        """Ingest movie data and yield a description."""
        temp_video_path = "temp_movie.mp4"
        with open(temp_video_path, "wb") as f:
            f.write(data)
        try:
            raw_frames, audio_file = self.process_video(temp_video_path)
            for frame in raw_frames:
                frame_text = process_frame_with_openai(
                    frame, self.openai_api_key
                )
                yield frame_text

            if audio_file:
                transcription_text = process_audio_with_openai(
                    audio_file, self.openai_api_key
                )
                yield transcription_text
        finally:
            os.remove(temp_video_path)

    def process_video(self, video_path):
        base64Frames = []
        base_video_path, _ = os.path.splitext(video_path)

        video = self.cv2.VideoCapture(video_path)
        total_frames = int(video.get(self.cv2.CAP_PROP_FRAME_COUNT))
        fps = video.get(self.cv2.CAP_PROP_FPS)
        frames_to_skip = int(fps * self.seconds_per_frame)
        curr_frame = 0

        # Calculate frames to skip based on max_frames if it is set
        if self.max_frames and self.max_frames < total_frames / frames_to_skip:
            frames_to_skip = max(total_frames // self.max_frames, 1)

        frame_count = 0
        while curr_frame < total_frames - 1 and (
            not self.max_frames or frame_count < self.max_frames
        ):
            video.set(self.cv2.CAP_PROP_POS_FRAMES, curr_frame)
            success, frame = video.read()
            if not success:
                break
            _, buffer = self.cv2.imencode(".jpg", frame)
            base64Frames.append(base64.b64encode(buffer).decode("utf-8"))
            curr_frame += frames_to_skip
            frame_count += 1
        video.release()

        audio_path = f"{base_video_path}.wav"
        audio_file = None
        with self.mp.VideoFileClip(video_path) as clip:
            if clip.audio is not None:
                clip.audio.write_audiofile(
                    audio_path, codec="pcm_s16le", fps=16000
                )
                audio_file = open(audio_path, "rb")
                os.remove(audio_path)

        return base64Frames, audio_file
